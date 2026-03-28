import streamlit as st
import requests
import PyPDF2
import docx
import json
import pytesseract
from PIL import Image
import pandas as pd
from pptx import Presentation

# --- CONFIGURATION ---
SOLR_URL = "http://localhost:8983/solr/docs"

# Chemin Tesseract (À adapter selon ton installation)
# pytesseract.pytesseract.tesseract_cmd = r'C:\Tesseract-OCR\tesseract.exe'

st.set_page_config(page_title="Solr Expert Search", layout="wide")

# --- STYLE CSS ---
st.markdown("""
    <style>
    mark { background-color: #ffeb3b; color: black; font-weight: bold; padding: 0 2px; border-radius: 3px; }
    div.stButton > button { border-radius: 20px; }
    .scroll-box {
        height: 400px; overflow-y: scroll; padding: 15px; border: 1px solid #444; 
        border-radius: 8px; background-color: #1e1e1e; color: #ffffff;
        white-space: pre-wrap; font-family: monospace;
    }
    </style>
""", unsafe_allow_html=True)

# --- FONCTIONS TECHNIQUES ---

def parse_facets(facet_list):
    """Transforme la liste Solr [val, count, val, count] en dictionnaire"""
    if not facet_list: return {}
    return {facet_list[i]: facet_list[i+1] for i in range(0, len(facet_list), 2)}

def extract_text(file):
    try:
        ext = file.name.split('.')[-1].lower()
        text_parts = []
        if ext in ["jpg", "jpeg", "png"]:
            text_parts.append(pytesseract.image_to_string(Image.open(file), lang='fra+eng'))
        elif ext == "pdf":
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                t = page.extract_text()
                if t: text_parts.append(t)
        elif ext == "docx":
            doc = docx.Document(file)
            for para in doc.paragraphs: text_parts.append(para.text)
        elif ext == "xlsx":
            text_parts.append(pd.read_excel(file).to_string(index=False))
        elif ext == "pptx":
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text_parts.append(shape.text)
        elif ext == "json":
            text_parts.append(json.dumps(json.load(file), indent=2, ensure_ascii=False))
        elif ext == "txt":
            text_parts.append(file.getvalue().decode("utf-8"))
        return "\n".join(text_parts).strip()
    except Exception as e:
        st.error(f"Erreur d'extraction : {e}")
        return ""

def get_solr_suggestions(prefix):
    if not prefix or len(prefix) < 2: return []
    try:
        r = requests.get(f"{SOLR_URL}/terms", params={"terms.fl": "content_t", "terms.prefix": prefix.lower(), "terms.limit": 8})
        terms = r.json().get("terms", {}).get("content_t", [])
        return [terms[i] for i in range(0, len(terms), 2)]
    except: return []

# --- INITIALISATION SESSION STATE ---
if "search_query" not in st.session_state: st.session_state.search_query = ""
if "ext_facets" not in st.session_state: st.session_state.ext_facets = {}
if "author_facets" not in st.session_state: st.session_state.author_facets = {}

# --- SIDEBAR : INDEXATION ---
with st.sidebar:
    st.header("📥 Indexation")
    up_file = st.file_uploader("Fichier", type=["pdf", "docx", "txt", "xlsx", "pptx", "json", "jpg", "png"])
    author = st.text_input("Auteur", value="Anonyme")
    keywords = st.text_input("Mots-clés")

    if up_file and st.button("🚀 Indexer"):
        content = extract_text(up_file)
        if content:
            data = [{
                "id": up_file.name, "content_t": content, "author_s": author,
                "keywords_t": keywords or "N/A", "extension_s": up_file.name.split('.')[-1].lower()
            }]
            requests.post(f"{SOLR_URL}/update?commit=true", json=data)
            st.success("Indexé !")
            st.rerun()

    if st.button("🗑️ VIDER L'INDEX"):
        requests.post(f"{SOLR_URL}/update?commit=true", json={"delete": {"query": "*:*"}})
        st.rerun()

# --- ZONE DE RECHERCHE ---
st.title("🔍 Recherche Documentaire Intelligente")

query_input = st.text_input("Rechercher...", value=st.session_state.search_query)

# Suggestions
if query_input:
    suggs = get_solr_suggestions(query_input)
    if suggs:
        cols = st.columns(len(suggs) + 1)
        cols[0].write("💡")
        for i, s in enumerate(suggs):
            if cols[i+1].button(s, key=f"s_{s}"):
                st.session_state.search_query = s
                st.rerun()

# --- ÉTAPE CRUCIALE : REQUÊTE SOLR AVANT L'AFFICHAGE DES FILTRES ---
res_docs = []
res_high = {}

if query_input:
    search_term = f"{query_input}*"
    params = {
        "q": f"content_t:{search_term} OR id:{search_term}",
        "facet": "true",
        "facet.field": ["extension_s", "author_s"],
        "facet.mincount": 1,
        "fl": "*,score",
        "hl": "true", "hl.fl": "content_t", "hl.fragsize": 0,
        "hl.simple.pre": "<mark>", "hl.simple.post": "</mark>", "hl.method": "unified"
    }
    try:
        r = requests.get(f"{SOLR_URL}/select", params=params)
        res = r.json()
        res_docs = res.get("response", {}).get("docs", [])
        res_high = res.get("highlighting", {})
        
        # On met à jour les facettes dans le session_state AVANT que les filtres ne soient dessinés
        facets_raw = res.get("facet_counts", {}).get("facet_fields", {})
        st.session_state.ext_facets = parse_facets(facets_raw.get("extension_s", []))
        st.session_state.author_facets = parse_facets(facets_raw.get("author_s", []))
    except: pass

# --- AFFICHAGE DES FILTRES (FACETTES) ---
c1, c2, c3 = st.columns(3)

with c1:
    options = ["pdf", "docx", "txt", "xlsx", "pptx", "json", "jpg", "png"]
    # Ici, st.session_state.ext_facets contient les chiffres de la requête juste au-dessus
    labels = {opt: f"{opt} ({st.session_state.ext_facets.get(opt, 0)})" for opt in options}
    f_types = st.multiselect("Extensions", options=options, format_func=lambda x: labels[x], default=options)

with c2:
    f_author = st.text_input("Auteur")
    if st.session_state.author_facets:
        st.caption("👤 " + ", ".join([f"{k}({v})" for k,v in st.session_state.author_facets.items()]))

with c3:
    target = st.radio("Cible", ["Tout", "Contenu", "Titre"], horizontal=True)

# --- AFFICHAGE DES RÉSULTATS ---
if query_input:
    if not res_docs:
        st.warning("Aucun résultat.")
    else:
        st.success(f"🎯 {len(res_docs)} document(s) trouvé(s)")
        for doc in res_docs:
            d_id = doc['id']
            # Filtrage manuel par extension (côté interface)
            if doc.get('extension_s') in f_types:
                with st.expander(f"📄 {d_id} (Score: {round(float(doc.get('score', 0)), 2)})"):
                    st.write(f"**Auteur :** {doc.get('author_s')} | **Mots-clés :** {doc.get('keywords_t')}")
                    content_high = res_high.get(d_id, {}).get('content_t', [doc.get('content_t', "")])[0]
                    st.markdown(f'<div class="scroll-box">{content_high}</div>', unsafe_allow_html=True)
                    st.download_button("📥 Texte brut", str(doc.get('content_t', "")), file_name=f"{d_id}.txt", key=f"dl_{d_id}")